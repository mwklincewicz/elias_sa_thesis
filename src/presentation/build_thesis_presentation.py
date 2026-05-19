from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT_DIR = ROOT / "output"
TEST_DATA_DIR = ROOT / "test data"
MODEL_DIR = OUTPUT_DIR / "bert_emotion_model"
EXTERNAL_DIR = TEST_DATA_DIR / "output" / "external_evaluation"
ASSET_DIR = OUTPUT_DIR / "presentation_assets"
DOWNLOADS_DIR = Path.home() / "Downloads"
DEFAULT_OUTPUT = DOWNLOADS_DIR / "lemotif_eda_and_training_results_presentation.pptx"

TITLE_COLOR = RGBColor(25, 43, 79)
ACCENT_COLOR = RGBColor(46, 117, 182)
LIGHT_BG = RGBColor(245, 247, 250)
TEXT_COLOR = RGBColor(40, 40, 40)
MUTED_COLOR = RGBColor(110, 110, 110)
MATPLOTLIB_TEXT_COLOR = "#282828"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def ensure_asset_dir() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def create_external_per_emotion_chart(metrics: dict, destination: Path) -> Path:
    label_metrics = pd.DataFrame(metrics["test_metrics"]["label_metrics"])
    label_metrics = label_metrics.sort_values("f1", ascending=True)

    sns.set_theme(style="whitegrid")
    fig, ax = plt.subplots(figsize=(10, 6))
    colors = sns.color_palette("crest", n_colors=len(label_metrics))
    ax.barh(label_metrics["label"], label_metrics["f1"], color=colors)
    ax.set_title("Stories External Test: Per-emotion F1", fontsize=16, pad=12)
    ax.set_xlabel("F1 score")
    ax.set_ylabel("")
    ax.set_xlim(0, max(0.6, float(label_metrics["f1"].max()) + 0.08))

    for index, row in label_metrics.iterrows():
        ax.text(
            float(row["f1"]) + 0.01,
            label_metrics.index.get_loc(index),
            f"n={int(row['support'])}",
            va="center",
            fontsize=9,
            color=MATPLOTLIB_TEXT_COLOR,
        )

    fig.tight_layout()
    fig.savefig(destination, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return destination


def create_training_comparison_chart(lemotif_metrics: dict, external_metrics: dict, destination: Path) -> Path:
    records = [
        {
            "metric": "Emotion micro F1",
            "Lemotif split": lemotif_metrics["test_metrics"]["micro"]["f1"],
            "Stories external": external_metrics["test_metrics"]["micro"]["f1"],
        },
        {
            "metric": "Emotion macro F1",
            "Lemotif split": lemotif_metrics["test_metrics"]["macro"]["f1"],
            "Stories external": external_metrics["test_metrics"]["macro"]["f1"],
        },
        {
            "metric": "Sentiment accuracy",
            "Lemotif split": lemotif_metrics["test_sentiment_metrics"]["accuracy"],
            "Stories external": external_metrics["test_sentiment_metrics"]["accuracy"],
        },
        {
            "metric": "Sentiment macro F1",
            "Lemotif split": lemotif_metrics["test_sentiment_metrics"]["macro"]["f1"],
            "Stories external": external_metrics["test_sentiment_metrics"]["macro"]["f1"],
        },
    ]

    frame = pd.DataFrame.from_records(records)
    melted = frame.melt(id_vars="metric", var_name="run", value_name="score")

    sns.set_theme(style="whitegrid")
    fig, ax = plt.subplots(figsize=(10, 5.6))
    sns.barplot(data=melted, x="metric", y="score", hue="run", palette=["#2F5597", "#5B9BD5"], ax=ax)
    ax.set_title("Training run comparison", fontsize=16, pad=12)
    ax.set_xlabel("")
    ax.set_ylabel("Score")
    ax.set_ylim(0, 1.0)
    ax.tick_params(axis="x", rotation=14)
    ax.legend(title="")

    for patch in ax.patches:
        height = patch.get_height()
        ax.annotate(
            f"{height:.3f}",
            (patch.get_x() + patch.get_width() / 2.0, height),
            ha="center",
            va="bottom",
            fontsize=9,
            xytext=(0, 4),
            textcoords="offset points",
        )

    fig.tight_layout()
    fig.savefig(destination, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return destination


def build_title_slide(prs: Presentation, lemotif_rows: int, stories_rows: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Emotion Recognition Thesis Results", "EDA across both datasets and outcomes from both BERT training runs")

    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.7), Inches(3.8))
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    title = frame.paragraphs[0]
    title.text = "Presentation flow"
    title.font.size = Pt(24)
    title.font.bold = True
    title.font.color.rgb = TITLE_COLOR

    for line in [
        f"Lemotif dataset EDA (n={lemotif_rows:,})",
        f"Stories external dataset EDA (n={stories_rows:,})",
        "In-domain BERT baseline results",
        "Full-Lemotif training with Stories external test results",
        "Closing comparison and thesis-ready takeaway",
    ]:
        paragraph = frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = TEXT_COLOR
        paragraph.space_before = Pt(8)

    note = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(5.5), Inches(0.8))
    note_frame = note.text_frame
    note_frame.text = "Planned length: about 5 to 7 minutes, with roughly 30 to 45 seconds per slide."
    note_frame.paragraphs[0].font.size = Pt(16)
    note_frame.paragraphs[0].font.color.rgb = MUTED_COLOR

    add_summary_card(
        slide,
        Inches(7.0),
        Inches(1.9),
        Inches(5.0),
        Inches(3.9),
        "Key message",
        [
            "Lemotif gives the stronger in-domain benchmark.",
            "Stories shows a meaningful domain-shift drop.",
            "Both datasets share the same label structure, so the comparison is interpretable.",
        ],
    )


def add_header(slide, title: str, subtitle: str) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.95))
    band.fill.solid()
    band.fill.fore_color.rgb = TITLE_COLOR
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.16), Inches(8.2), Inches(0.35))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.font.size = Pt(26)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(0.57), Inches(0.53), Inches(8.5), Inches(0.22))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(11)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(227, 233, 245)


def add_summary_card(slide, left, top, width, height, heading: str, lines: Iterable[str]) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = ACCENT_COLOR
    shape.line.width = Pt(1.5)

    box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.15), width - Inches(0.36), height - Inches(0.3))
    frame = box.text_frame
    frame.word_wrap = True
    p0 = frame.paragraphs[0]
    p0.text = heading
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = TITLE_COLOR

    for line in lines:
        paragraph = frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = TEXT_COLOR
        paragraph.space_before = Pt(6)


def add_footer_note(slide, text: str) -> None:
    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(6.8), Inches(12.4), Inches(0.42))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(236, 242, 248)
    note.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.62), Inches(6.88), Inches(12.0), Inches(0.2))
    frame = box.text_frame
    frame.text = text
    frame.paragraphs[0].font.size = Pt(11)
    frame.paragraphs[0].font.color.rgb = MUTED_COLOR
    frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def add_picture_with_caption(slide, image_path: Path, caption: str, left, top, width, height) -> None:
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    cap_box = slide.shapes.add_textbox(left, top + height + Inches(0.03), width, Inches(0.22))
    cap_frame = cap_box.text_frame
    cap_frame.text = caption
    cap_frame.paragraphs[0].font.size = Pt(10.5)
    cap_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    cap_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def build_four_figure_slide(prs: Presentation, title: str, subtitle: str, items: list[tuple[Path, str]], footer: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, subtitle)

    left_positions = [Inches(0.55), Inches(6.65)]
    top_positions = [Inches(1.15), Inches(4.0)]
    width = Inches(5.9)
    height = Inches(2.3)

    for index, (image_path, caption) in enumerate(items):
        row = index // 2
        col = index % 2
        add_picture_with_caption(slide, image_path, caption, left_positions[col], top_positions[row], width, height)

    add_footer_note(slide, footer)


def build_three_figure_slide(prs: Presentation, title: str, subtitle: str, items: list[tuple[Path, str]], footer: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, subtitle)

    left_positions = [Inches(0.45), Inches(4.45), Inches(8.45)]
    width = Inches(4.0)
    height = Inches(5.2)
    top = Inches(1.18)

    for index, (image_path, caption) in enumerate(items):
        add_picture_with_caption(slide, image_path, caption, left_positions[index], top, width, height)

    add_footer_note(slide, footer)


def build_lemotif_results_slide(prs: Presentation, metrics: dict, model_figure_dir: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Training 1: Lemotif Split Baseline", "BERT fine-tuned on the standard train/validation/test split")

    add_summary_card(
        slide,
        Inches(0.55),
        Inches(1.18),
        Inches(3.35),
        Inches(5.25),
        "Key results",
        [
            f"Split sizes: {metrics['split_sizes']['train']} train / {metrics['split_sizes']['val']} val / {metrics['split_sizes']['test']} test",
            f"Emotion micro F1: {metrics['test_metrics']['micro']['f1']:.3f}",
            f"Emotion macro F1: {metrics['test_metrics']['macro']['f1']:.3f}",
            f"Sentiment accuracy: {metrics['test_sentiment_metrics']['accuracy']:.3f}",
            f"Threshold selected on validation: {metrics['best_threshold']:.2f}",
            "Interpretation: strong first benchmark, but still recall-heavy and much better on frequent labels.",
        ],
    )

    add_picture_with_caption(
        slide,
        model_figure_dir / "fig1_training_progress.png",
        "Training and validation curves across four epochs",
        Inches(4.15),
        Inches(1.2),
        Inches(4.2),
        Inches(2.35),
    )
    add_picture_with_caption(
        slide,
        model_figure_dir / "fig2_validation_vs_test_metrics.png",
        "Validation and test metrics stay fairly aligned",
        Inches(8.55),
        Inches(1.2),
        Inches(4.2),
        Inches(2.35),
    )
    add_picture_with_caption(
        slide,
        model_figure_dir / "fig3_per_emotion_f1.png",
        "Per-emotion test performance highlights strong labels and rare-label instability",
        Inches(4.15),
        Inches(3.95),
        Inches(4.2),
        Inches(2.35),
    )
    add_picture_with_caption(
        slide,
        model_figure_dir / "fig6_sentiment_confusion_matrix.png",
        "Derived sentiment is much easier than full emotion recognition",
        Inches(8.55),
        Inches(3.95),
        Inches(4.2),
        Inches(2.35),
    )

    add_footer_note(slide, "Suggested talk track: the model learns the common positive emotions best, while exact multi-label matching remains very difficult.")


def build_external_results_slide(
    prs: Presentation,
    lemotif_metrics: dict,
    external_metrics: dict,
    comparison_chart: Path,
    external_per_emotion_chart: Path,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Training 2: Full Lemotif Train, Stories External Test", "Same label space, but evaluated on a new dataset with visible domain shift")

    add_summary_card(
        slide,
        Inches(0.55),
        Inches(1.18),
        Inches(3.35),
        Inches(5.25),
        "External-test results",
        [
            f"Training rows: {external_metrics['split_sizes']['train']} | Stories test rows: {external_metrics['split_sizes']['test']}",
            f"Emotion micro F1: {external_metrics['test_metrics']['micro']['f1']:.3f}",
            f"Emotion macro F1: {external_metrics['test_metrics']['macro']['f1']:.3f}",
            f"Sentiment accuracy: {external_metrics['test_sentiment_metrics']['accuracy']:.3f}",
            f"Avg predicted labels per entry: {external_metrics['test_metrics']['avg_predicted_labels_per_entry']:.2f}",
            "Interpretation: performance drops sharply out of domain, especially for positive sentiment and rarer emotions.",
        ],
    )

    add_picture_with_caption(
        slide,
        comparison_chart,
        "Direct comparison of the two training runs",
        Inches(4.15),
        Inches(1.2),
        Inches(4.15),
        Inches(2.35),
    )
    add_picture_with_caption(
        slide,
        external_per_emotion_chart,
        "Stories external per-emotion F1",
        Inches(8.5),
        Inches(1.2),
        Inches(4.25),
        Inches(2.35),
    )

    takeaway = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.15), Inches(3.95), Inches(8.6), Inches(2.35))
    takeaway.fill.solid()
    takeaway.fill.fore_color.rgb = RGBColor(250, 252, 255)
    takeaway.line.color.rgb = ACCENT_COLOR

    box = slide.shapes.add_textbox(Inches(4.35), Inches(4.12), Inches(8.2), Inches(1.95))
    frame = box.text_frame
    p0 = frame.paragraphs[0]
    p0.text = "Comparison takeaway"
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = TITLE_COLOR

    drops = [
        f"Emotion micro F1 falls from {lemotif_metrics['test_metrics']['micro']['f1']:.3f} to {external_metrics['test_metrics']['micro']['f1']:.3f}.",
        f"Emotion macro F1 falls from {lemotif_metrics['test_metrics']['macro']['f1']:.3f} to {external_metrics['test_metrics']['macro']['f1']:.3f}.",
        f"Sentiment accuracy falls from {lemotif_metrics['test_sentiment_metrics']['accuracy']:.3f} to {external_metrics['test_sentiment_metrics']['accuracy']:.3f}.",
        "This is a useful thesis result: the baseline generalizes only partially when the writing style and label distribution change.",
    ]

    for line in drops:
        paragraph = frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = TEXT_COLOR
        paragraph.space_before = Pt(6)

    add_footer_note(slide, "Suggested talk track: present the Stories evaluation dataset as an external validity check rather than as a failure case.")


def build_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Closing Takeaways", "Short ending slide for the last 30 to 45 seconds")

    add_summary_card(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(5.7),
        Inches(4.9),
        "What the results show",
        [
            "Both datasets share the same emotion/topic structure, so the comparison is coherent.",
            "Lemotif EDA confirms a highly imbalanced, multi-label emotion task.",
            "The in-domain BERT baseline is usable as a first benchmark.",
            "External testing on Stories reveals a clear domain-shift penalty.",
            "This supports a thesis conclusion that better robustness needs more diverse data and possibly stronger modeling.",
        ],
    )

    add_summary_card(
        slide,
        Inches(6.85),
        Inches(1.45),
        Inches(5.7),
        Inches(4.9),
        "Suggested final sentence",
        [
            "The model captures the main emotional signal within Lemotif,",
            "but performance drops on the external Stories set,",
            "which suggests that future work should focus on domain adaptation,",
            "richer supervision, and stronger external validation.",
        ],
    )

    add_footer_note(slide, "End on the contrast between a solid in-domain benchmark and a weaker but very informative external test.")


def build_presentation(destination: Path) -> Path:
    ensure_asset_dir()
    lemotif_metrics = load_json(MODEL_DIR / "metrics.json")
    external_metrics = load_json(EXTERNAL_DIR / "metrics.json")

    lemotif_rows = len(pd.read_csv(OUTPUT_DIR / "lemotif_cleaned.csv"))
    stories_rows = len(pd.read_csv(TEST_DATA_DIR / "stories_data.csv"))

    comparison_chart = create_training_comparison_chart(
        lemotif_metrics,
        external_metrics,
        ASSET_DIR / "training_run_comparison.png",
    )
    external_per_emotion_chart = create_external_per_emotion_chart(
        external_metrics,
        ASSET_DIR / "stories_external_per_emotion_f1.png",
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title_slide(prs, lemotif_rows=lemotif_rows, stories_rows=stories_rows)

    lemotif_eda_groups = [
        (
            "Lemotif EDA (1/3)",
            "Core distribution view of the cleaned Lemotif dataset",
            [
                (OUTPUT_DIR / "fig1_word_count_distribution.png", "Figure 1. Word count distribution"),
                (OUTPUT_DIR / "fig2_emotion_prevalence.png", "Figure 2. Emotion prevalence"),
                (OUTPUT_DIR / "fig3_topic_prevalence.png", "Figure 3. Topic prevalence"),
                (OUTPUT_DIR / "fig4_emotion_corr.png", "Figure 4. Emotion correlation matrix"),
            ],
            "Talk track: start with label imbalance and the fact that reflections often cluster around a few common positive emotions and topics.",
        ),
        (
            "Lemotif EDA (2/3)",
            "Relationship plots and representative examples",
            [
                (OUTPUT_DIR / "fig5_topic_emotion.png", "Figure 5. Topic-emotion heatmap"),
                (OUTPUT_DIR / "fig6_emotion_labels_per_entry.png", "Figure 6. Emotion labels per entry"),
                (OUTPUT_DIR / "fig7_topic_labels_per_entry.png", "Figure 7. Topic labels per entry"),
                (OUTPUT_DIR / "fig8_representative_entries.png", "Figure 8. Representative entries"),
            ],
            "Talk track: emphasize that Lemotif is genuinely multi-label, which makes exact prediction harder than simple single-label classification.",
        ),
    ]

    for title, subtitle, items, footer in lemotif_eda_groups:
        build_four_figure_slide(prs, title, subtitle, items, footer)

    build_three_figure_slide(
        prs,
        "Lemotif EDA (3/3)",
        "Network and conditional-probability views",
        [
            (OUTPUT_DIR / "fig9_emotion_cooccurrence_network.png", "Figure 9. Emotion co-occurrence network"),
            (OUTPUT_DIR / "fig10_length_vs_emotion_probability.png", "Figure 10. Length versus emotion probability"),
            (OUTPUT_DIR / "fig11_topic_to_emotion_probability.png", "Figure 11. Topic to emotion probability"),
        ],
        "Talk track: the final three figures show co-occurrence structure and how text length or topic context shifts emotion likelihoods.",
    )

    stories_eda_groups = [
        (
            "Stories EDA (1/3)",
            "External dataset prepared into the same structure as Lemotif",
            [
                (TEST_DATA_DIR / "output" / "eda" / "fig1_word_count_distribution.png", "Figure 1. Word count distribution"),
                (TEST_DATA_DIR / "output" / "eda" / "fig2_emotion_prevalence.png", "Figure 2. Emotion prevalence"),
                (TEST_DATA_DIR / "output" / "eda" / "fig3_topic_prevalence.png", "Figure 3. Topic prevalence"),
                (TEST_DATA_DIR / "output" / "eda" / "fig4_emotion_corr.png", "Figure 4. Emotion correlation matrix"),
            ],
            "Talk track: introduce Stories as comparable in structure, but noticeably different in tone, topic balance, and emotion mix.",
        ),
        (
            "Stories EDA (2/3)",
            "Topic-emotion structure and labeling density in the external set",
            [
                (TEST_DATA_DIR / "output" / "eda" / "fig5_topic_emotion.png", "Figure 5. Topic-emotion heatmap"),
                (TEST_DATA_DIR / "output" / "eda" / "fig6_emotion_labels_per_entry.png", "Figure 6. Emotion labels per entry"),
                (TEST_DATA_DIR / "output" / "eda" / "fig7_topic_labels_per_entry.png", "Figure 7. Topic labels per entry"),
                (TEST_DATA_DIR / "output" / "eda" / "fig8_representative_entries.png", "Figure 8. Representative entries"),
            ],
            "Talk track: note where Stories resembles Lemotif and where it diverges, because those divergences help explain the external-test drop later.",
        ),
    ]

    for title, subtitle, items, footer in stories_eda_groups:
        build_four_figure_slide(prs, title, subtitle, items, footer)

    build_three_figure_slide(
        prs,
        "Stories EDA (3/3)",
        "Co-occurrence and conditional-probability patterns in the external set",
        [
            (TEST_DATA_DIR / "output" / "eda" / "fig9_emotion_cooccurrence_network.png", "Figure 9. Emotion co-occurrence network"),
            (TEST_DATA_DIR / "output" / "eda" / "fig10_length_vs_emotion_probability.png", "Figure 10. Length versus emotion probability"),
            (TEST_DATA_DIR / "output" / "eda" / "fig11_topic_to_emotion_probability.png", "Figure 11. Topic to emotion probability"),
        ],
        "Talk track: finish the Stories section by connecting these patterns to why the model may generalize only partially across datasets.",
    )

    build_lemotif_results_slide(prs, lemotif_metrics, MODEL_DIR / "figures")
    build_external_results_slide(prs, lemotif_metrics, external_metrics, comparison_chart, external_per_emotion_chart)
    build_closing_slide(prs)

    destination.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(destination))
    return destination


def main() -> None:
    output_path = build_presentation(DEFAULT_OUTPUT)
    print(f"Saved presentation to: {output_path}")


if __name__ == "__main__":
    main()
